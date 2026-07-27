import sys
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ---------------------------------------------------------
    # Color Palette Definitions (Premium Dark Mode)
    # ---------------------------------------------------------
    BG_COLOR = RGBColor(11, 15, 23)        # Deep Navy (#0B0F17)
    CARD_BG = RGBColor(23, 29, 41)         # Card Dark Gray (#171D29)
    ACCENT_TEAL = RGBColor(0, 212, 170)    # Vibrant Teal (#00D4AA)
    ACCENT_BLUE = RGBColor(59, 130, 246)   # Blue (#3B82F6)
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(138, 153, 173)   # Gray/Blue (#8A99AD)
    
    # Helper to add dark background to a slide
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper to add standard slide header
    def add_slide_header(slide, title_text, category_text="PHARMACI PLATFORM PORTFOLIO"):
        # Category indicator
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Plus Jakarta Sans"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_TEAL
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Plus Jakarta Sans"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    blank_layout = prs.slide_layouts[6]

    # ---------------------------------------------------------
    # SLIDE 1: Title Slide (The Strategic Hook)
    # ---------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    accent_strip = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15)
    )
    accent_strip.fill.solid()
    accent_strip.fill.fore_color.rgb = ACCENT_TEAL
    accent_strip.line.fill.background()
    
    main_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tf1 = main_box.text_frame
    tf1.word_wrap = True
    
    p_main_title = tf1.paragraphs[0]
    p_main_title.text = "PharmaCI Copilot"
    p_main_title.font.name = "Plus Jakarta Sans"
    p_main_title.font.size = Pt(64)
    p_main_title.font.bold = True
    p_main_title.font.color.rgb = TEXT_WHITE
    p_main_title.space_after = Pt(15)
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "AI-Powered Competitive Intelligence for Specialty Pharma"
    p_sub.font.name = "Plus Jakarta Sans"
    p_sub.font.size = Pt(24)
    p_sub.font.bold = True
    p_sub.font.color.rgb = ACCENT_TEAL
    p_sub.space_after = Pt(10)
    
    p_sub2 = tf1.add_paragraph()
    p_sub2.text = "A Credible GTM Path to $5M ARR via Product-Led Growth & Multi-Step Agentic CI"
    p_sub2.font.name = "Plus Jakarta Sans"
    p_sub2.font.size = Pt(16)
    p_sub2.font.color.rgb = TEXT_MUTED
    p_sub2.space_after = Pt(40)
    
    p_author = tf1.add_paragraph()
    p_author.text = "Alok Pandey  |  Candidate: AI Technical Product Owner — Pharma AI Platforms"
    p_author.font.name = "Plus Jakarta Sans"
    p_author.font.size = Pt(14)
    p_author.font.bold = True
    p_author.font.color.rgb = TEXT_WHITE

    # ---------------------------------------------------------
    # SLIDE 2: The Core Market Problem (Underserved Mid-Tier)
    # ---------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_slide_header(slide2, "The Underserved Mid-Tier: Legacy Cost vs. Resource Scarcity")
    
    # Column 1 Card (The Enterprise Giants)
    col1_bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    col1_bg.fill.solid()
    col1_bg.fill.fore_color.rgb = CARD_BG
    col1_bg.line.fill.background()
    
    col1_text = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_c1 = col1_text.text_frame
    tf_c1.word_wrap = True
    
    p = tf_c1.paragraphs[0]
    p.text = "THE ENTERPRISE REALITY"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(14)
    
    bullets_c1 = [
        "Top-10 global pharmaceutical giants spend millions annually on legacy platforms (IQVIA, Veeva, EvaluatePharma).",
        "They employ dedicated competitive intelligence (CI) and commercial analyst teams to clean and query data.",
        "Brand strategies are refreshed slowly due to complex extraction requests."
    ]
    for b in bullets_c1:
        bp = tf_c1.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)

    # Column 2 Card (The Specialty / Mid-Tier Pain)
    col2_bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    col2_bg.fill.solid()
    col2_bg.fill.fore_color.rgb = CARD_BG
    col2_bg.line.fill.background()
    
    col2_text = slide2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_c2 = col2_text.text_frame
    tf_c2.word_wrap = True
    
    p = tf_c2.paragraphs[0]
    p.text = "THE MID-TIER SPECIALTY GAP"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    bullets_c2 = [
        "Mid-tier specialty brand teams ($100M - $1B revenues) have zero dedicated competitive intelligence analysts.",
        "Brand Leads and Market Access Directors waste 15+ hours a week manually parsing trials, payer filings, and news feed tables.",
        "Crucial pricing actions and pipeline changes are missed due to information overload."
    ]
    for b in bullets_c2:
        bp = tf_c2.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)

    # ---------------------------------------------------------
    # SLIDE 3: The Product Wedge (Core Pillars)
    # ---------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_slide_header(slide3, "The Product Wedge: From Simple RAG to Autonomous Agents")
    
    pillars = [
        ("1. Autonomous ReAct Agent Loop", "Uses Groq llama-3.3-70b-versatile. Decides which database tools to call dynamically (queries, trends, cross-references) instead of returning static keyword hits.", ACCENT_TEAL),
        ("2. Interactive Metrics & Drill-downs", "Allows brand leads to click any high-level card metric on the dashboard to immediately view corresponding underlying signals in an inspection panel.", TEXT_WHITE),
        ("3. Compliance Audit Diagnostics", "Renders Ingestion authorities, dynamic posture flags, and grounding metrics. Proves absolute auditability to conservative medical and legal stakeholders.", TEXT_WHITE),
        ("4. Fail-Safe Rate Limit Interceptor", "Auto-detects rate limits (429 errors) and transitions to local trace simulations, guaranteeing zero interruption for the review panel.", ACCENT_TEAL)
    ]
    
    for i, (title, desc, color) in enumerate(pillars):
        row = i // 2
        col = i % 2
        
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.8 + row * 2.5)
        
        card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.6), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()
        
        card_text = slide3.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.2), Inches(1.8))
        tf = card_text.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Plus Jakarta Sans"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(8)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Plus Jakarta Sans"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 4: FEATURE WALKTHROUGH - Strategic Analytics Dashboard
    # ---------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_slide_header(slide4, "Feature 1: Interactive Strategic Analytics Dashboard", "WALKTHROUGH: TAB 1")
    
    # Left Content Column
    left_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "EXECUTIVE HEALTH & METRICS MATRIX"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    f1_bullets = [
        "Share of Voice (SOV) Index: Dynamic donut chart tracking competitor signal intensity across channels.",
        "Competitor Posture Mix: Classifies signals into Offensive, Defensive, and Neutral postures using Plotly visual layouts.",
        "Intensity Heatmap Matrix: Cross-references Competitor vs. Signal Type, highlighting areas of competitor concentration.",
        "Interactivity: Directly integrates with sidebar filters to update market charts in under 50ms."
    ]
    for b in f1_bullets:
        bp = tf_l.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)

    # Right Content Column (Visual Card Mockup)
    right_bg = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(1.8), Inches(5.0), Inches(4.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = CARD_BG
    right_bg.line.fill.background()
    
    right_text = slide4.shapes.add_textbox(Inches(7.6), Inches(2.2), Inches(4.6), Inches(4.0))
    tf_r = right_text.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "KEY METRIC CARDS"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(20)
    
    metrics = [
        ("Signals Traced", "21 Active Signals"),
        ("High Impact Triggers", "6 Action Required"),
        ("Offensive Competitor Moves", "6 Outward Pressures")
    ]
    for name, stat in metrics:
        bp = tf_r.add_paragraph()
        bp.text = f"{name}: "
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.bold = True
        bp.font.color.rgb = TEXT_WHITE
        
        run = bp.add_run()
        run.text = stat
        run.font.bold = True
        run.font.color.rgb = ACCENT_TEAL
        bp.space_after = Pt(14)

    # ---------------------------------------------------------
    # SLIDE 5: FEATURE WALKTHROUGH - Competitive Monitoring Feed
    # ---------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_slide_header(slide5, "Feature 2: Competitive Feed & Dynamic Implication Analysis", "WALKTHROUGH: TAB 2")
    
    left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "INTELLIGENT FEED DESTRUCTURING"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(14)
    
    f2_bullets = [
        "Impact Badging: Color-coded severity indicators (High/Medium/Low) based on clinical and regulatory urgency.",
        "Dynamic Implication Analysis: Groq API evaluates the 'So-what' implication for brand leads in under 25 words.",
        "Auditable Metadata: Every signal card is backed by source tracking (Payer Filings, Field Intelligence, etc.).",
        "Strategic Segmentation: Highlights competitor intent (Offensive/Defensive) for quick scanning."
    ]
    for b in f2_bullets:
        bp = tf_l.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)

    # Right Content Box
    right_bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(1.8), Inches(5.0), Inches(4.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = CARD_BG
    right_bg.line.fill.background()
    
    right_text = slide5.shapes.add_textbox(Inches(7.6), Inches(2.2), Inches(4.6), Inches(4.0))
    tf_r = right_text.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "AUDIT TRACE DETAILS"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(20)
    
    audit_points = [
        ("Source Ingestion", "E.g., Payer Filings, News outlets, Clinical Registries"),
        ("Diagnostic Panel", "Tracks request ingestion authority and pipeline performance metrics."),
        ("Bypasses / Safety", "Protects client-facing systems from arbitrary text extrapolation.")
    ]
    for t_t, t_d in audit_points:
        bp = tf_r.add_paragraph()
        bp.text = f"{t_t}: "
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(13)
        bp.font.bold = True
        bp.font.color.rgb = TEXT_WHITE
        
        run = bp.add_run()
        run.text = t_d
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        bp.space_after = Pt(14)

    # ---------------------------------------------------------
    # SLIDE 6: FEATURE WALKTHROUGH - Grounded Q&A RAG
    # ---------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_slide_header(slide6, "Feature 3: Grounded Q&A Retrieval (RAG)", "WALKTHROUGH: TAB 3")
    
    left_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "GROUNDED KNOWLEDGE RETRIEVAL"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    f3_bullets = [
        "Keyword Search Ingestion: Mapped to parse assets, companies, and therapeutic keywords locally in under 1ms.",
        "Response Synthesis: Groq compiles matching signals into a direct, 100-word brief for brand teams.",
        "Factuality Lock: Strictly instructed to state data limitations rather than speculating, satisfying pharma regulatory checks.",
        "Preset Queries: Ready-to-run buttons populate search strings instantly to show functionality."
    ]
    for b in f3_bullets:
        bp = tf_l.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)

    # Right Content Box
    right_bg = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(1.8), Inches(5.0), Inches(4.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = CARD_BG
    right_bg.line.fill.background()
    
    right_text = slide6.shapes.add_textbox(Inches(7.6), Inches(2.2), Inches(4.6), Inches(4.0))
    tf_r = right_text.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "EVIDENCE TRANSPARENCY"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(20)
    
    p_ev = tf_r.add_paragraph()
    p_ev.text = "Source Grounding:"
    p_ev.font.name = "Plus Jakarta Sans"
    p_ev.font.size = Pt(14)
    p_ev.font.bold = True
    p_ev.font.color.rgb = TEXT_WHITE
    p_ev.space_after = Pt(8)
    
    p_ev_desc = tf_r.add_paragraph()
    p_ev_desc.text = "Every RAG output displays the exact database records (dates, headlines, drugs) cited. No output is shown without explicit link references, building user confidence."
    p_ev_desc.font.name = "Plus Jakarta Sans"
    p_ev_desc.font.size = Pt(13)
    p_ev_desc.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 7: FEATURE WALKTHROUGH - Autonomous Agent Analyst Loop
    # ---------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_slide_header(slide7, "Feature 4: Autonomous Agent Analyst Loop (ReAct)", "WALKTHROUGH: TAB 4")
    
    left_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "MULTISTEP REASONING WORKFLOW"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    f4_bullets = [
        "Multi-Step Planning: The agent decides dynamically which tools to execute in what sequence to answer complex queries.",
        "Interactive Status Trace: Employs Streamlit st.status to show step spinners and parameters live in the UI.",
        "Tool Execution: Decouples operations into query_signals, get_trend, cross_reference, and draft_recommendation.",
        "Guardrail Meter: Calculates database coverage to verify grounding integrity before final synthesis."
    ]
    for b in f4_bullets:
        bp = tf_l.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)

    # Right Content Box
    right_bg = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(1.8), Inches(5.0), Inches(4.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = CARD_BG
    right_bg.line.fill.background()
    
    right_text = slide7.shapes.add_textbox(Inches(7.6), Inches(2.2), Inches(4.6), Inches(4.0))
    tf_r = right_text.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "THE REASONING TRACE"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(20)
    
    trace_steps = [
        ("Step 1: query_signals", "Filters GSK respiratory updates."),
        ("Step 2: cross_reference", "Maps Sun Pharma field rep metrics."),
        ("Step 3: draft_recommendation", "Drafts commercial posture changes.")
    ]
    for step, desc in trace_steps:
        bp = tf_r.add_paragraph()
        bp.text = f"{step}: "
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(13)
        bp.font.bold = True
        bp.font.color.rgb = ACCENT_TEAL
        
        run = bp.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(10)

    # ---------------------------------------------------------
    # SLIDE 8: FEATURE WALKTHROUGH - Executive Briefing Digests
    # ---------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_slide_header(slide8, "Feature 5: Compiled Executive Digests & Export", "WALKTHROUGH: TAB 5")
    
    left_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "EXECUTIVE SUMMARY COMPILATION"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(14)
    
    f5_bullets = [
        "Digest Compilation: Aggregates current filtered competitor update cards into one briefing document with one click.",
        "Inline Recommendations: Appends dynamic commercial implications directly under signal summaries.",
        "One-Click Download: Exports compile digests as Markdown files ready to copy into slides or executive emails.",
        "Formatting Structure: Segmented by Competitor and Drug focus for clean executive delivery."
    ]
    for b in f5_bullets:
        bp = tf_l.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)

    # Right Content Box
    right_bg = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(1.8), Inches(5.0), Inches(4.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = CARD_BG
    right_bg.line.fill.background()
    
    right_text = slide8.shapes.add_textbox(Inches(7.6), Inches(2.2), Inches(4.6), Inches(4.0))
    tf_r = right_text.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "GTM UTILITY"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(20)
    
    p_ut = tf_r.add_paragraph()
    p_ut.text = "The 'Work-Saver' Wedge:"
    p_ut.font.name = "Plus Jakarta Sans"
    p_ut.font.size = Pt(14)
    p_ut.font.bold = True
    p_ut.font.color.rgb = TEXT_WHITE
    p_ut.space_after = Pt(8)
    
    p_ut_desc = tf_r.add_paragraph()
    p_ut_desc.text = "In enterprise sales, users convert when a tool makes their day-to-day admin tasks disappear. The Briefing compilation tool allows brand managers to build their weekly alignment reports instantly, securing high product retention."
    p_ut_desc.font.name = "Plus Jakarta Sans"
    p_ut_desc.font.size = Pt(13)
    p_ut_desc.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 9: Target Segment & Ideal Customer Profile (ICP)
    # ---------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_slide_header(slide9, "Target Market: High-Growth Specialty Pharma Portfolios", "COMMERCIAL STRATEGY")
    
    left_box = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "IDEAL CUSTOMER PROFILE (ICP)"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    points = [
        ("Company Type", "Mid-tier specialty pharmaceutical developers ($100M - $1.5B revenue) in the US and Europe."),
        ("Target Domains", "Complex therapeutic areas with high competitor activity: Oncology, Dermatology, Immunology, Diabetes/Obesity, and Vaccines."),
        ("The Buyer Personas", "Brand Commercial Directors, Market Access Leads, and Medical Affairs Directors."),
        ("Speed Advantage", "While legacy platforms require 6+ months of setup, PharmaCI can ingest and index local datasets in 15 minutes.")
    ]
    for t_title, t_desc in points:
        bp = tf_l.add_paragraph()
        bp.text = f"•  {t_title}: "
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.bold = True
        bp.font.color.rgb = TEXT_WHITE
        
        run = bp.add_run()
        run.text = t_desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        bp.space_after = Pt(12)
        
    right_bg = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(1.8), Inches(5.0), Inches(4.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = CARD_BG
    right_bg.line.fill.background()
    
    right_text = slide9.shapes.add_textbox(Inches(7.6), Inches(2.2), Inches(4.6), Inches(4.0))
    tf_r = right_text.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "THE REVENUE OPPORTUNITY"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(20)
    
    p_num = tf_r.add_paragraph()
    p_num.text = "350+"
    p_num.font.name = "Plus Jakarta Sans"
    p_num.font.size = Pt(56)
    p_num.font.bold = True
    p_num.font.color.rgb = ACCENT_TEAL
    p_num.space_after = Pt(5)
    
    p_lbl = tf_r.add_paragraph()
    p_lbl.text = "Mid-tier specialty pharma companies in US & Europe representing a $140M+ Addressable TA Module Market. Capturing just 50 modules achieves our $5M ARR target."
    p_lbl.font.name = "Plus Jakarta Sans"
    p_lbl.font.size = Pt(14)
    p_lbl.font.color.rgb = TEXT_WHITE

    # ---------------------------------------------------------
    # SLIDE 10: The Product-Led Growth (PLG) Sales Loop
    # ---------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_slide_header(slide10, "GTM Playbook: 3-Step Product-Led Growth (PLG) Loop", "COMMERCIAL STRATEGY")
    
    steps = [
        ("Step 1: Custom Data Hook", "Before the first meeting, we ingest the prospect's actual therapeutic competitor data into the dashboard. No generic templates are used.", ACCENT_BLUE),
        ("Step 2: The 'Wow' Demo", "We let the prospect enter a complex query. The Reasoning Trace spins live in front of them—showing the agent planning, fetching, and citing sources.", ACCENT_TEAL),
        ("Step 3: Sandbox Conversion", "We hand over a 14-day trial. VPs use the 'Compile Digest' tab to auto-generate markdown briefings for weekly leadership meetings, driving organic user conversion.", ACCENT_BLUE)
    ]
    
    for i, (title, desc, color) in enumerate(steps):
        left = Inches(0.8 + i * 4.0)
        top = Inches(1.8)
        
        card = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()
        
        bar = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.7), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        card_text = slide10.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), Inches(3.3), Inches(4.2))
        tf = card_text.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Plus Jakarta Sans"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(14)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Plus Jakarta Sans"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 11: The Math to $5M ARR (Commercial Model)
    # ---------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_slide_header(slide11, "The Path to $5M ARR: Revenue Model & Projections", "COMMERCIAL STRATEGY")
    
    left_box = slide11.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "THE LAND & EXPAND MATH"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    math_bullets = [
        "Unit Price: $100,000 / year per Therapeutic Area (TA) Module.",
        "Target Conversion: 50 closed modules globally (less than 3.6% market penetration).",
        "Logo Scaling: We only need 10 enterprise accounts buying 5 modules each to reach this target.",
        "Expansion Potential: Specialty pharma accounts expand budget into adjacent brand units within 6 months of initial onboarding."
    ]
    for b in math_bullets:
        bp = tf_l.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)
        
    right_bg = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(1.8), Inches(5.0), Inches(4.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = CARD_BG
    right_bg.line.fill.background()
    
    tbl_text = slide11.shapes.add_textbox(Inches(7.6), Inches(2.0), Inches(4.6), Inches(4.4))
    tf_t = tbl_text.text_frame
    tf_t.word_wrap = True
    
    p = tf_t.paragraphs[0]
    p.text = "YEAR 1 PHASED PROJECTIONS"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(20)
    
    quarters = [
        ("Q1 Launch & Pilot", "5 Beta Pilots Closed", "$250K"),
        ("Q2 Expansion", "10 New Modules Closed", "$1,000K"),
        ("Q3 Land & Expand", "15 Modules + Cross-sell", "$1,500K"),
        ("Q4 Scale Run-rate", "20 Modules Closed", "$2,250K")
    ]
    for q_name, q_action, q_rev in quarters:
        qp = tf_t.add_paragraph()
        qp.text = f"{q_name}: "
        qp.font.name = "Plus Jakarta Sans"
        qp.font.size = Pt(13)
        qp.font.bold = True
        qp.font.color.rgb = TEXT_WHITE
        
        run1 = qp.add_run()
        run1.text = f"{q_action} — "
        run1.font.bold = False
        run1.font.color.rgb = TEXT_MUTED
        
        run2 = qp.add_run()
        run2.text = q_rev
        run2.font.bold = True
        run2.font.color.rgb = ACCENT_TEAL
        
        qp.space_after = Pt(14)

    # ---------------------------------------------------------
    # SLIDE 12: Technical Gaps & Future Scale Plan (Roadmap)
    # ---------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_slide_header(slide12, "Product Roadmap: Transitioning to Enterprise Scale", "TECHNICAL ROADMAP")
    
    sc1 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    sc1.fill.solid()
    sc1.fill.fore_color.rgb = CARD_BG
    sc1.line.fill.background()
    
    sc1_text = slide12.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_sc1 = sc1_text.text_frame
    tf_sc1.word_wrap = True
    
    p = tf_sc1.paragraphs[0]
    p.text = "SHORT-TERM UPGRADES (Q1 - Q2)"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)
    
    st_bullets = [
        "Hybrid Search Ingestion: Combine dense vector embeddings with client-side BM25 keyword matching for clinical PDF structures.",
        "Multi-Agent Orchestration: Transition from single-agent ReAct loops to dedicated, graph-based agents (LangGraph) for clinical, regulatory, and market access analysis.",
        "Source Reliability Weighting: Score database updates by ingestion authority."
    ]
    for b in st_bullets:
        bp = tf_sc1.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)

    sc2 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    sc2.fill.solid()
    sc2.fill.fore_color.rgb = CARD_BG
    sc2.line.fill.background()
    
    sc2_text = slide12.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_sc2 = sc2_text.text_frame
    tf_sc2.word_wrap = True
    
    p = tf_sc2.paragraphs[0]
    p.text = "LONG-TERM SCALING (Q3 - Q4)"
    p.font.name = "Plus Jakarta Sans"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(14)
    
    lt_bullets = [
        "Continuous Data Integration: Connect pipelines to live API portals (ClinicalTrials.gov, FDA registers, Payer Formularies).",
        "LLM Observability Suites: Setup LangSmith / Arize Phoenix to monitor trace token cost, request latency, and tool selection metrics.",
        "Governance & Defense: Add LLM input sanitization and guardrail libraries (Llama Guard) to shield client data."
    ]
    for b in lt_bullets:
        bp = tf_sc2.add_paragraph()
        bp.text = "•  " + b
        bp.font.name = "Plus Jakarta Sans"
        bp.font.size = Pt(14)
        bp.font.color.rgb = TEXT_WHITE
        bp.space_after = Pt(12)

    prs.save(r"C:\Users\aloks\Downloads\PharmaCI_GTM_Pitch_Deck.pptx")
    print("Presentation created successfully at C:\\Users\\aloks\\Downloads\\PharmaCI_GTM_Pitch_Deck.pptx")

if __name__ == "__main__":
    create_presentation()
